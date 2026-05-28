"""
Report engine — orchestrates generation of all campaign deliverables.

Every campaign produces Markdown + JSON reports plus Maltego CSV export.
Reports are saved to ./reports/<client>/<engagement_id>/<timestamp>/.
"""

from __future__ import annotations

import json
from datetime import datetime
from pathlib import Path
from typing import Any

import re

import structlog

log = structlog.get_logger(__name__)


def strip_agent_scaffolding(text: str) -> str:
    """Remove machine protocol scaffolding from agent prose before it lands
    in a human report (Wave F-B6).

    The agent executor emits ``FINDINGS_JSON:[...]`` for the findings
    pipeline; when an agent's raw output is rendered verbatim, that marker
    and its JSON array leak into the deliverable (the executive summary in
    the 2026-05-27 run carried a literal ``FINDINGS_JSON:[{...}]`` blob).
    This keeps only the human-readable prose: everything from the
    ``FINDINGS_JSON:`` marker to the end of its JSON array is stripped.
    """
    if not text:
        return text
    marker = "FINDINGS_JSON:"
    idx = text.find(marker)
    while idx != -1:
        after = text[idx + len(marker):].lstrip()
        rest_offset = len(text) - len(after)
        try:
            _parsed, consumed = json.JSONDecoder().raw_decode(after)
            end = rest_offset + consumed
        except (json.JSONDecodeError, ValueError):
            # Can't parse the array; drop to end of the line as a fallback.
            nl = text.find("\n", idx)
            end = len(text) if nl == -1 else nl
        text = (text[:idx] + text[end:])
        idx = text.find(marker)
    # Collapse the blank gap the removal may leave behind.
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _provider_has_evidence(data: Any) -> bool:
    """True when a cloud_intel entry carries a positive signal, not just the
    fact that a recon tool ran (Wave F-B4).

    aws_recon / gcp_recon write a state entry even when they find nothing,
    so keying "provider detected" on entry existence fabricates presence.
    Require a non-empty evidence field instead.
    """
    if not isinstance(data, dict):
        return False
    evidence_keys = (
        "openid_config", "user_realm", "onmicrosoft_domain", "s3_buckets",
        "buckets", "public_buckets", "projects", "instances", "functions",
        "account_id", "accounts", "services", "tenant_id",
    )
    for k in evidence_keys:
        v = data.get(k)
        if isinstance(v, (list, dict, tuple)) and v:
            return True
        if isinstance(v, str) and v and v.lower() != "unknown":
            return True
    return False


def _code_source_has_evidence(data: Any) -> bool:
    """True when a code_intel entry actually surfaced something (leaks,
    findings, repos, dependencies) rather than merely being queried (F-B4)."""
    if not isinstance(data, dict):
        return False
    for k in ("leaks", "findings", "repos", "repositories", "secrets"):
        if data.get(k):
            return True
    inner = data.get("data")
    if isinstance(inner, dict):
        for k in ("dependencies", "packages", "repos", "results"):
            if inner.get(k):
                return True
    return False


class ReportEngine:
    """
    Generates all campaign deliverables from the accumulated state.

    Reports:
    1. Executive summary (1 page)
    2. Full engagement report
    3. Asset inventory (JSON + Markdown + CSV)
    4. Phishing target package
    5. Cloud & identity posture brief
    6. Attack surface matrix
    7. Maltego CSV export
    8. Full findings JSON
    9. Campaign metadata JSON
    """

    def __init__(self, campaign_id: str, engagement_id: str, scope_hash: str, output_dir: Path) -> None:
        self.campaign_id = campaign_id
        self.engagement_id = engagement_id
        self.scope_hash = scope_hash
        # Pin the version of the framework that produced the report.
        # Pairs with ``scope_hash`` in every footer so an auditor can
        # reproduce the run: same scope file + same code = same report
        # (modulo upstream provider drift). Resolved lazily ── if the
        # package import is somehow broken at runtime, the report
        # still ships with ``unknown`` rather than crashing the
        # campaign at deliverable time.
        try:
            from nexusrecon import __version__ as _pkg_version
            self.nexusrecon_version: str = str(_pkg_version)
        except Exception:
            self.nexusrecon_version = "unknown"
        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.report_paths: dict[str, str] = {}

    # ── Main orchestration ─────────────────────────────────────────────────────

    def generate_all(self, state: dict[str, Any]) -> dict[str, str]:
        """Generate all report types. Returns dict of report_name → file_path."""
        log.info("Generating all reports", campaign_id=self.campaign_id)

        self.report_paths["top_threads"] = self._top_threads_to_pull(state)
        self.report_paths["executive_summary"] = self._executive_summary(state)
        self.report_paths["full_report"] = self._full_report(state)
        self.report_paths["asset_inventory"] = self._asset_inventory(state)
        self.report_paths["phishing_package"] = self._phishing_package(state)
        self.report_paths["cloud_posture"] = self._cloud_posture(state)
        self.report_paths["attack_surface"] = self._attack_surface(state)
        self.report_paths["findings_json"] = self._findings_json(state)
        self.report_paths["campaign_meta"] = self._campaign_meta(state)
        self.report_paths["people_map"] = self._people_map(state)
        self.report_paths["vuln_correlation"] = self._vuln_correlation(state)
        self.report_paths["vendor_supply_chain"] = self._vendor_supply_chain(state)
        self.report_paths["jira_tracker"] = self._jira_tracker(state)
        self.report_paths["entity_graph_html"] = self._entity_graph_html(state)
        self.report_paths["pdf_report"] = self._pdf_report(state)
        self.report_paths["pptx_report"] = self._pptx_report(state)
        self.report_paths["harvested_credentials"] = self._harvested_credentials(state)
        # D7: credential exposure paths — personal pivot + punch list
        self.report_paths["credential_exposure_paths"] = self._credential_exposure_paths(state)
        # E11: spear-phishing intelligence — per-target pretext dossiers
        md_path, json_path = self._spear_phishing_intelligence(state)
        self.report_paths["spear_phishing_intelligence"] = md_path
        self.report_paths["pretext_candidates_json"] = json_path
        # V3 Move 2: master_report runs LAST so it can link to every other report
        self.report_paths["master_report"] = self._master_report(state)

        # Phase 1 of the toolchain plan: parallel Obsidian-flavored
        # master report. Gated on the --obsidian CLI flag (lives in
        # state["generate_obsidian"]). Must run AFTER _master_report
        # because it reads the standard file from disk and transforms
        # it — no source-of-truth fork.
        if state.get("generate_obsidian"):
            self.report_paths["master_report_obsidian"] = (
                self._master_report_obsidian(state)
            )

        return self.report_paths

    # ── Top Threads to Pull ────────────────────────────────────────────────────

    def _top_threads_to_pull(self, state: dict[str, Any]) -> str:
        """
        Operator-facing 'where to start' document — the single most important
        deliverable for kicking off the rest of the pentest campaign.
        """
        ranked_threads = state.get("ranked_threads", [])
        state.get("findings", [])
        state.get("vuln_intel", {}).get("enriched_cves", {})
        state.get("email_intel", {}).get("emails", {})

        lines = [
            "# Top 10 Threads to Pull",
            "",
            f"**Campaign:** {self.campaign_id}",
            f"**Engagement:** {self.engagement_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            "",
            "> This document is your starting point. Each thread is a specific, actionable",
            "> attack path ranked by the probability it leads to a successful compromise.",
            "> Work top-to-bottom.",
            "",
            "---",
            "",
        ]

        if not ranked_threads:
            lines.append("*No ranked threads available — ensure Phase 7 and Phase 8 completed successfully.*")
            lines.extend(self._render_coverage_section(state))
            path = self.output_dir / "top_threads.md"
            path.write_text("\n".join(lines), encoding="utf-8")
            return str(path)

        for i, thread in enumerate(ranked_threads[:10], 1):
            score_pct = f"{thread.get('score', 0) * 100:.0f}%"
            sev = thread.get("severity", "info").upper()
            title = thread.get("title", "Untitled")
            category = thread.get("category", "general")
            confidence = thread.get("confidence", 0.0)

            lines.extend([
                f"## Thread {i}: {title}",
                "",
                "| Field | Value |",
                "|-------|-------|",
                f"| **Priority Score** | {score_pct} |",
                f"| **Severity** | {sev} |",
                f"| **Category** | {category.upper()} |",
                f"| **Confidence** | {confidence:.0%} |",
            ])

            # CVE-specific metadata
            if thread.get("cve_id"):
                lines.append(f"| **CVE** | {thread['cve_id']} |")
            if thread.get("cvss"):
                lines.append(f"| **CVSS** | {thread['cvss']} |")
            if thread.get("epss") is not None:
                lines.append(f"| **EPSS** | {thread['epss']:.1%} probability of exploitation in 30 days |")

            lines.append("")

            # Flags
            flags = []
            if thread.get("in_kev"):
                flags.append("🔴 **CISA KEV** — confirmed actively exploited in the wild")
            if thread.get("has_metasploit"):
                flags.append("🔴 **Metasploit module available** — weaponized, low skill required")
            elif thread.get("has_exploit"):
                flags.append("🟠 **Public PoC exploit** — requires adaptation")
            if thread.get("has_nuclei_template"):
                flags.append("🟢 **Nuclei template available** — verify with one command")
            if thread.get("cloud_provider"):
                flags.append(f"☁️ **Cloud ({thread['cloud_provider'].upper()})** — no auth required to access")

            if flags:
                lines.append("**Risk Indicators:**")
                lines.append("")
                for flag in flags:
                    lines.append(f"- {flag}")
                lines.append("")

            lines.append(f"**Description:** {thread.get('description', '')[:400]}")
            lines.append("")

            if thread.get("affected_assets"):
                assets_str = ", ".join(f"`{a}`" for a in thread["affected_assets"][:5])
                lines.append(f"**Affected Assets:** {assets_str}")
                lines.append("")

            steps = thread.get("next_steps", [])
            if steps:
                lines.append("**Recommended Actions (in order):**")
                lines.append("")
                for j, step in enumerate(steps, 1):
                    lines.append(f"{j}. {step}")
                lines.append("")

            if thread.get("breach_sources"):
                lines.append(f"**Breach Sources:** {', '.join(thread['breach_sources'][:5])}")
                lines.append("")

            lines.append("---")
            lines.append("")

        # Summary table
        lines.extend([
            "## Summary Matrix",
            "",
            "| # | Score | Sev | Category | Title | KEV | MSF | Template |",
            "|---|-------|-----|----------|-------|-----|-----|---------|",
        ])
        for i, t in enumerate(ranked_threads[:10], 1):
            sev = t.get("severity", "?")[0].upper()
            score_pct = f"{t.get('score', 0) * 100:.0f}%"
            cat = t.get("category", "?")[:8]
            title = (t.get("title") or "")[:40]
            kev = "✓" if t.get("in_kev") else ""
            msf = "✓" if t.get("has_metasploit") else ""
            tpl = "✓" if t.get("has_nuclei_template") else ""
            lines.append(f"| {i} | {score_pct} | {sev} | {cat} | {title} | {kev} | {msf} | {tpl} |")

        lines.extend(self._render_coverage_section(state))

        lines.extend([
            "",
            "---",
            "",
            "*Generated by NexusRecon scoring engine. Scores are relative within this campaign.*",
        ])

        path = self.output_dir / "top_threads.md"
        path.write_text("\n".join(lines), encoding="utf-8")

        # Surface validated harvested credentials as standalone threads
        harvested_creds = state.get("harvested_credentials", [])
        for cred in harvested_creds:
            score = 0.95 if cred.get("validated") else 0.7
            ranked_threads.append({
                "title": f"Harvested credential: {cred.get('cred_type', 'unknown')} from {cred.get('source_type', '?')}",
                "score": score,
                "severity": "critical" if cred.get("validated") else "high",
                "category": "harvested_credential",
                "confidence": score,
                "description": f"Credential ({cred.get('cred_type')}) found at {cred.get('source_url', '?')}. "
                               f"Validated: {cred.get('validated', False)}. Hash: {cred.get('value_hash', '')[:16]}...",
                "next_steps": cred.get("next_steps", []),
            })

        # Also write JSON for programmatic use
        threads_json_path = self.output_dir / "top_threads.json"
        threads_json_path.write_text(
            json.dumps({"campaign_id": self.campaign_id, "threads": ranked_threads[:10]}, indent=2, default=str),
            encoding="utf-8",
        )

        return str(path)

    def _render_coverage_section(self, state: dict[str, Any]) -> list[str]:
        """Coverage / what-we-checked appendix (Wave F-B1).

        Renders the absence-of-evidence and below-floor items the scoring
        engine split out of the ranked threads, plus (when the run-health
        summary is available) tools that were degraded or failed, labelled
        as "not assessed" rather than a clean negative. The goal is an
        honest record of scope that does NOT create follow-up work: a
        pentester opening top_threads sees attack surface first, with the
        dead ends quarantined here.
        """
        coverage = state.get("coverage_items", []) or []
        run_health = state.get("run_health", {}) or {}
        degraded = run_health.get("degraded", []) or []
        if not coverage and not degraded:
            return []

        out = [
            "",
            "---",
            "",
            "## Coverage / What We Checked",
            "",
            "Items below are NOT attack surface. They record where we looked and",
            "found nothing actionable, or where a tool could not complete. Kept for",
            "an honest scope record, not as follow-up work.",
            "",
        ]
        if degraded:
            out.append("**Not assessed (tool failed or returned implausibly empty):**")
            out.append("")
            for d in degraded:
                out.append(f"- `{d.get('tool', '?')}`: {d.get('reason', '')}")
            out.append("")
        if coverage:
            out.append("**Checked, nothing actionable found:**")
            out.append("")
            for c in coverage:
                conf = c.get("confidence", 0.0) or 0.0
                out.append(f"- {c.get('title', '')} _(confidence {conf:.0%})_")
            out.append("")
        return out

    # ── Executive Summary ──────────────────────────────────────────────────────

    def _executive_summary(self, state: dict[str, Any]) -> str:
        """1-page red-team focused summary."""
        findings = state.get("findings", [])
        critical = [f for f in findings if f.get("severity") == "critical"]
        high = [f for f in findings if f.get("severity") == "high"]
        medium = [f for f in findings if f.get("severity") == "medium"]

        subdomains = state.get("subdomain_intel", {})
        emails = state.get("email_intel", {}).get("emails", {})
        cloud_intel = state.get("cloud_intel", {})

        lines = [
            "# Executive Summary — NexusRecon Campaign",
            "",
            f"**Campaign ID:** {self.campaign_id}",
            f"**Engagement ID:** {self.engagement_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            f"**Scope Hash:** {self.scope_hash}",
            f"**Tooling:** NexusRecon v{self.nexusrecon_version}",
            "",
            "---",
            "",
            "## Overview",
            "",
            "This report summarizes the reconnaissance findings for the authorized engagement.",
            "",
            f"- **Total Findings:** {len(findings)}",
            f"- **Critical:** {len(critical)}",
            f"- **High:** {len(high)}",
            f"- **Medium:** {len(medium)}",
            f"- **Subdomains Discovered:** {len(subdomains)}",
            f"- **Email Addresses:** {len(emails)}",
            # B35: count verified cloud presence comprehensively, not just S3 buckets.
            # Counts each cloud provider with attribution_confidence >= 0.5 plus
            # bucket/storage objects regardless of provider.
            f"- **Cloud Assets:** {_count_cloud_assets(cloud_intel)}",
            "",
            "## Key Findings",
            "",
        ]

        for i, f in enumerate(findings[:10], 1):
            lines.append(f"{i}. **[{f.get('severity', 'info').upper()}]** {f.get('title', 'Untitled')}")
            lines.append(f"   - {f.get('description', '')[:200]}")
            lines.append(f"   - Source: {f.get('source', 'unknown')} | Confidence: {f.get('confidence', 0):.0%}")
            lines.append("")

        # Top threads to pull
        ranked_threads = state.get("ranked_threads", [])
        if ranked_threads:
            lines.extend([
                "## Top Threads to Pull",
                "",
                "Ranked by exploitability × impact × exposure. Start here.",
                "",
            ])
            for i, thread in enumerate(ranked_threads[:10], 1):
                score_pct = f"{thread.get('score', 0) * 100:.0f}%"
                sev = thread.get("severity", "info").upper()
                title = thread.get("title", "")
                lines.append(f"### {i}. [{sev}] {title} — Priority Score: {score_pct}")
                lines.append("")
                lines.append(f"{thread.get('description', '')[:250]}")
                lines.append("")
                flags = []
                if thread.get("in_kev"):
                    flags.append("KEV (actively exploited in the wild)")
                if thread.get("has_metasploit"):
                    flags.append("Metasploit module available")
                elif thread.get("has_exploit"):
                    flags.append("Public PoC exploit exists")
                if thread.get("has_nuclei_template"):
                    flags.append("Nuclei template available — runnable immediately")
                if flags:
                    lines.append(f"**Flags:** {' | '.join(flags)}")
                    lines.append("")
                steps = thread.get("next_steps", [])
                if steps:
                    lines.append("**Immediate next steps:**")
                    for step in steps[:3]:
                        lines.append(f"1. {step}")
                    lines.append("")

        # Risk analyst agent synthesis
        analyst_msg = next(
            (m for m in reversed(state.get("agent_messages", [])) if m.get("agent") == "risk_analyst"),
            None,
        )
        if analyst_msg and analyst_msg.get("analysis"):
            assessment = strip_agent_scaffolding(analyst_msg["analysis"])
            if assessment:
                lines.extend([
                    "## Analyst Assessment",
                    "",
                    assessment,
                    "",
                ])

        lines.extend([
            "---",
            "",
            "*Full details in the complete engagement report.*",
        ])

        path = self.output_dir / "executive_summary.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)

    # ── Full Report ────────────────────────────────────────────────────────────

    def _full_report(self, state: dict[str, Any]) -> str:
        """Complete engagement report with methodology."""
        findings = state.get("findings", [])
        completed = state.get("completed_phases", [])

        lines = [
            "# NexusRecon Engagement Report",
            "",
            f"**Campaign ID:** {self.campaign_id}",
            f"**Engagement ID:** {self.engagement_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            f"**Scope Hash:** {self.scope_hash}",
            f"**Tooling:** NexusRecon v{self.nexusrecon_version}",
            "",
            "---",
            "",
            "## Table of Contents",
            "",
            "1. Executive Summary",
            "2. Methodology",
            "3. Findings",
            "4. Infrastructure Analysis",
            "5. Cloud & Identity Analysis",
            "6. Code & Secret Exposure",
            "7. Attack Surface Matrix",
            "8. Recommendations",
            "9. Appendix",
            "",
            "---",
            "",
            "## 1. Executive Summary",
            "",
            "See executive_summary.md for the one-page summary.",
            "",
            "## 2. Methodology",
            "",
            "This engagement followed a phased reconnaissance methodology:",
            "",
        ]

        phases = {
            "phase1": "Passive Footprinting (T0) — domain, DNS, subdomain, certificate enumeration",
            "phase2": "Identity & Cloud (T0-T1) — M365/AWS/GCP enumeration, email harvesting",
            "phase3": "Deep Subdomain & Code (T0) — recursive subdomain, GitHub/code scanning",
            "phase4": "Correlation & Hypothesis — cross-source validation, lead generation",
            "phase5": "Light Active (T2) — HTTP probing, screenshots, tech fingerprinting",
            "phase6": "Active (T3) — brute force, content fuzzing (if authorized)",
            "phase7": "Pretext & Vulnerability Correlation — social engineering intel, CVE mapping",
            "phase8": "Attack Surface Prioritization — scoring, ranking, PRE-ATT&CK mapping",
            "phase9": "Reporting — deliverable generation",
        }

        for p in completed:
            lines.append(f"- **{p}:** {phases.get(p, '')}")

        lines.extend([
            "",
            "## 3. Findings",
            "",
        ])

        for i, f in enumerate(findings, 1):
            lines.extend([
                f"### {i}. {f.get('title', 'Untitled')}",
                "",
                f"- **Severity:** {f.get('severity', 'info').upper()}",
                f"- **Confidence:** {f.get('confidence', 0):.0%}",
                f"- **Category:** {f.get('category', 'unknown')}",
                f"- **Source:** {f.get('source', 'unknown')}",
                f"- **Timestamp:** {f.get('timestamp', 'unknown')}",
                f"- **Evidence Hash:** {f.get('raw_evidence_hash', 'N/A')}",
                "",
                f"**Description:** {f.get('description', '')}",
                "",
            ])
            # B27: conditional rendering — skip empty/N/A optional fields entirely
            assets = [a for a in (f.get("affected_assets") or []) if a]
            if assets:
                lines.extend([f"**Affected Assets:** {', '.join(str(a) for a in assets[:5])}", ""])
            mitre = [m for m in (f.get("mitre_techniques") or []) if m]
            if mitre:
                lines.extend([f"**MITRE Techniques:** {', '.join(str(m) for m in mitre)}", ""])
            rec = (f.get("recommendation") or "").strip()
            if rec and rec.lower() not in ("n/a", "-", "none"):
                lines.extend([f"**Recommendation:** {rec}", ""])
            lines.extend(["---", ""])

        lines.extend([
            "## 4-9. Detailed Sections",
            "",
            "See individual report files for infrastructure, cloud, code, attack surface, and recommendation details.",
        ])

        path = self.output_dir / "full_report.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)

    # ── Asset Inventory ────────────────────────────────────────────────────────

    def _asset_inventory(self, state: dict[str, Any]) -> str:
        """Complete asset list in Markdown + JSON + CSV."""
        subdomains = state.get("subdomain_intel", {})
        emails = state.get("email_intel", {}).get("emails", {})
        cloud_intel = state.get("cloud_intel", {})

        # Markdown
        lines = [
            "# Asset Inventory",
            "",
            f"**Campaign:** {self.campaign_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            "",
            f"## Subdomains ({len(subdomains)})",
            "",
        ]
        for sub in sorted(subdomains.keys())[:500]:
            lines.append(f"- `{sub}`")

        lines.extend([
            "",
            f"## Email Addresses ({len(emails)})",
            "",
        ])
        for em, info in sorted(emails.items()):
            lines.append(f"- `{em}` — source: {info.get('source', 'unknown')}")

        lines.extend([
            "",
            "## Cloud Assets",
            "",
        ])
        for key, data in cloud_intel.items():
            lines.append(f"### {key}")
            for bucket in data.get("s3_buckets", []):
                lines.append(f"- S3: `{bucket.get('name', '')}` (public: {bucket.get('public', False)})")

        path_md = self.output_dir / "asset_inventory.md"
        path_md.write_text("\n".join(lines), encoding="utf-8")

        # JSON
        path_json = self.output_dir / "asset_inventory.json"
        inventory = {
            "campaign_id": self.campaign_id,
            "scope_hash": self.scope_hash,
            "nexusrecon_version": self.nexusrecon_version,
            "generated": datetime.utcnow().isoformat(),
            "subdomains": sorted(subdomains.keys()),
            "emails": list(emails.keys()),
            "cloud_assets": cloud_intel,
        }
        path_json.write_text(json.dumps(inventory, indent=2, default=str), encoding="utf-8")

        # CSV
        import csv
        path_csv = self.output_dir / "asset_inventory.csv"
        with open(path_csv, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Type", "Value", "Source", "Details"])
            for sub in sorted(subdomains.keys()):
                writer.writerow(["subdomain", sub, subdomains[sub].get("sources", ["unknown"])[0], ""])
            for em, info in sorted(emails.items()):
                writer.writerow(["email", em, info.get("source", "unknown"), info.get("position", "")])

        return str(path_md)

    # ── Phishing Package ───────────────────────────────────────────────────────

    def _phishing_package(self, state: dict[str, Any]) -> str:
        """Emails + pretext hooks + per-employee bundles + DMARC analysis."""
        from nexusrecon.core.identity_hygiene import is_probable_test_identity
        # F-B5: never build a pretext bundle for an obviously synthetic/test
        # address (abcfoo@, noreply@). They are not people; targeting them is
        # wasted operator effort and pollutes the package.
        emails = {
            em: info
            for em, info in state.get("email_intel", {}).get("emails", {}).items()
            if not is_probable_test_identity(em)
        }
        _raw_format = state.get("email_intel", {}).get("format", {})
        # Handle both legacy string format ("first.last@example.com") and
        # rich dict format ({"most_likely_pattern": ..., "most_likely_confidence": ...})
        if isinstance(_raw_format, str):
            email_format: dict[str, Any] = {
                "most_likely_pattern": _raw_format,
                "most_likely_confidence": 1.0,
            }
        else:
            email_format = _raw_format if isinstance(_raw_format, dict) else {}
        domain_intel = state.get("domain_intel", {})
        email_sec = domain_intel.get("dns", {})
        breach_intel = state.get("breach_intel", {})
        ranked_threads = state.get("ranked_threads", [])

        # Collect breach-linked and infostealer-linked emails for cross-reference
        breached_emails: set[str] = set()
        infostealer_emails: set[str] = set()
        for em, info in emails.items():
            if isinstance(info, dict):
                if info.get("breaches") or breach_intel.get(em):
                    breached_emails.add(em)
                if info.get("stealer_logs") or info.get("infostealer_hits"):
                    infostealer_emails.add(em)

        # Executive targets
        exec_keywords = ["ceo", "cfo", "cto", "ciso", "cso", "vp", "director", "executive", "president", "founder"]
        exec_targets: list[str] = []
        for em, info in emails.items():
            if isinstance(info, dict):
                pos = str(info.get("position", "")).lower()
                dept = str(info.get("department", "")).lower()
                if any(kw in pos or kw in dept for kw in exec_keywords):
                    exec_targets.append(em)

        lines = [
            "# Phishing Target Package",
            "",
            f"**Campaign:** {self.campaign_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            "",
            "## Email Format Analysis",
            "",
            f"Most likely pattern: `{email_format.get('most_likely_pattern', 'unknown')}` "
            f"(confidence: {email_format.get('most_likely_confidence', 0):.0%})",
            "",
        ]

        if email_format.get("pattern_distribution"):
            lines.append("### Pattern Distribution")
            for pattern, data in email_format["pattern_distribution"].items():
                lines.append(f"- `{pattern}`: {data.get('count', 0)} occurrences ({data.get('confidence', 0):.0%})")
            lines.append("")

        # ── Attack surface context ───────────────────────────────────────────
        breach_thread = next((t for t in ranked_threads if t.get("category") == "breach"), None)
        if breach_thread:
            lines.extend([
                "## ⚠ Breach Intelligence",
                "",
                f"**{breach_thread.get('title', '')}**",
                "",
                breach_thread.get("description", "")[:400],
                "",
            ])
            for step in breach_thread.get("next_steps", [])[:3]:
                lines.append(f"- {step}")
            lines.append("")

        # ── Email Security Posture ────────────────────────────────────────────
        dmarc = email_sec.get("dmarc_record", ["not found"])
        dmarc_val = dmarc[0] if isinstance(dmarc, list) and dmarc else str(dmarc)
        spf = email_sec.get("spf_record", ["not found"])
        spf_val = spf[0] if isinstance(spf, list) and spf else str(spf)

        lines.extend([
            "## Email Security Posture",
            "",
            "| Control | Status |",
            "|---------|--------|",
            f"| DMARC | `{dmarc_val[:120]}` |",
            f"| SPF | `{spf_val[:120]}` |",
            "",
        ])

        # DMARC policy assessment
        dmarc_lower = dmarc_val.lower()
        if "p=reject" in dmarc_lower:
            lines.append("**DMARC:** p=reject — direct spoofing of the root domain is blocked. Use subdomain or lookalike domains.")
        elif "p=quarantine" in dmarc_lower:
            lines.append("**DMARC:** p=quarantine — spoofed emails land in spam. Lookalike domains are more effective.")
        elif "p=none" in dmarc_lower or "not found" in dmarc_lower:
            lines.append("**DMARC:** p=none or absent — direct domain spoofing is viable. Craft emails from the exact target domain.")
        lines.append("")

        # ── Per-employee pretext bundles ─────────────────────────────────────
        lines.extend([
            "## Per-Employee Pretext Bundles",
            "",
            "Prioritised by: executive role > infostealer hit > breach hit > role.",
            "",
        ])

        # Build priority-sorted target list
        def _target_priority(em: str) -> int:
            score = 0
            if em in infostealer_emails:
                score += 40
            if em in exec_targets:
                score += 30
            if em in breached_emails:
                score += 20
            info = emails.get(em, {})
            if isinstance(info, dict):
                pos = info.get("position", "")
                if pos and pos not in ("unknown", "None", ""):
                    score += 10
            return score

        sorted_targets = sorted(emails.keys(), key=_target_priority, reverse=True)

        for em in sorted_targets[:50]:
            info = emails.get(em, {}) if isinstance(emails.get(em), dict) else {}
            role = info.get("position") or "Employee"
            dept = info.get("department") or "Unknown dept."
            source = info.get("source", "harvested")

            flags_str = ""
            tag_parts = []
            if em in infostealer_emails:
                tag_parts.append("INFOSTEALER")
            if em in exec_targets:
                tag_parts.append("EXECUTIVE")
            if em in breached_emails:
                tag_parts.append("BREACHED")
            flags_str = " | ".join(tag_parts)

            lines.extend([
                f"### `{em}`",
                "",
                f"- **Role:** {role}",
                f"- **Department:** {dept}",
                f"- **Source:** {source}",
            ])
            if flags_str:
                lines.append(f"- **Tags:** {flags_str}")

            # Build pretext suggestions based on role + breach status
            pretexts = _generate_pretext_hooks(em, role, dept, em in breached_emails, em in infostealer_emails)
            if pretexts:
                lines.append("")
                lines.append("**Pretext Hooks:**")
                lines.append("")
                for pt in pretexts:
                    lines.append(f"- {pt}")

            lines.append("")

        # ── Executive target summary ──────────────────────────────────────────
        if exec_targets:
            lines.extend([
                "## High-Value Targets (Executive)",
                "",
                "| Email | Role | Breached | Infostealer |",
                "|-------|------|----------|-------------|",
            ])
            for em in exec_targets[:20]:
                info = emails.get(em, {}) if isinstance(emails.get(em), dict) else {}
                role = info.get("position", "N/A")
                breached = "✓" if em in breached_emails else ""
                stealer = "✓" if em in infostealer_emails else ""
                lines.append(f"| `{em}` | {role} | {breached} | {stealer} |")
            lines.append("")

        # Phishing draft generator (opt-in via generate_phishing_drafts flag)
        if state.get("generate_phishing_drafts"):
            try:
                import asyncio as _asyncio

                from nexusrecon.core.config import get_config
                from nexusrecon.graph.agent_executor import AgentExecutor
                from nexusrecon.reports.phishing_drafts import generate_phishing_drafts
                executor = AgentExecutor(get_config())
                drafts = _asyncio.run(generate_phishing_drafts(state, executor, self.output_dir, max_targets=10))
                lines.append("\n## Generated Drafts\n")
                for target, draft_path in drafts.items():
                    if not target.startswith("__"):
                        lines.append(f"- [{target}]({Path(draft_path).name})")
            except Exception as exc:
                log.warning("Phishing draft generation failed", error=str(exc))

        path = self.output_dir / "phishing_package.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)

    # ── Cloud Posture ──────────────────────────────────────────────────────────

    def _cloud_posture(self, state: dict[str, Any]) -> str:
        """M365 federation, AWS account, and public cloud asset summary."""
        cloud_intel = state.get("cloud_intel", {})

        lines = [
            "# Cloud & Identity Posture Brief",
            "",
            f"**Campaign:** {self.campaign_id}",
            "",
        ]

        # Build each provider subsection from POSITIVE signal only, and emit
        # the header only when there is something to show (F-B4). The old
        # renderer printed a "## aws" header with "S3 Buckets Found: 0" and
        # an empty "## gcp" section purely because the recon tool ran; that
        # is noise, and "Tenant ID: unknown" reads like a real value.
        any_section = False
        for key, data in cloud_intel.items():
            if not isinstance(data, dict):
                continue
            section: list[str] = []

            oc = data.get("openid_config") or {}
            tid = oc.get("tenant_id")
            if tid and str(tid).lower() != "unknown":
                section.append(f"- Tenant ID: {tid}")
            iss = oc.get("issuer")
            if iss and str(iss).lower() != "unknown":
                section.append(f"- Issuer: {iss}")

            ur = data.get("user_realm")
            if isinstance(ur, dict) and ur:
                section.append(
                    f"- Federation: {'Federated (ADFS)' if ur.get('is_federated') else 'Managed'}"
                )

            acct = data.get("account_id")
            if acct and str(acct).lower() != "unknown":
                section.append(f"- Account ID: {acct}")
            accounts = data.get("accounts") or []
            if accounts:
                section.append("- Accounts: " + ", ".join(str(a) for a in accounts))
            services = data.get("services") or []
            if services:
                section.append("- Services: " + ", ".join(str(s) for s in services))

            s3 = data.get("s3_buckets") or []
            if s3:
                section.append(f"- S3 Buckets Found: {len(s3)}")
                for b in s3:
                    section.append(f"  - `{b.get('name', '')}` (public: {b.get('public', False)})")

            onm = data.get("onmicrosoft_domain") or {}
            if isinstance(onm, dict) and onm.get("domains"):
                domain_strs = [
                    str(d.get("domain") or d.get("tenant_id") or "?")
                    for d in onm["domains"]
                    if isinstance(d, dict)
                ]
                if domain_strs:
                    section.append("- onmicrosoft.com domains: " + ", ".join(domain_strs))

            if section:
                any_section = True
                lines.append(f"## {key}")
                lines.append("")
                lines.extend(section)
                lines.append("")

        if not any_section:
            lines.append(
                "_No positive cloud or identity posture signal was collected. The "
                "cloud recon tools ran but returned nothing actionable; see the "
                "run health summary for tool status._"
            )
            lines.append("")

        path = self.output_dir / "cloud_posture.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)

    # ── Attack Surface ─────────────────────────────────────────────────────────

    def _attack_surface(self, state: dict[str, Any]) -> str:
        """Likelihood × impact matrix with PRE-ATT&CK mapping."""
        from nexusrecon.core.scoring import likelihood_impact
        findings = state.get("findings", [])

        lines = [
            "# Attack Surface Matrix",
            "",
            f"**Campaign:** {self.campaign_id}",
            "",
            "| # | Severity | Confidence | Title | PRE-ATT&CK | Likelihood | Impact |",
            "|---|----------|------------|-------|-------------|------------|--------|",
        ]

        for i, f in enumerate(findings, 1):
            mitre = ", ".join(f.get("mitre_techniques", [])) or "-"
            # F-B6: compute Likelihood/Impact once in core/scoring.py and
            # render the real numbers, instead of the old blank "- | -" that
            # let the LLM exec-summary invent its own inconsistent integers.
            likelihood, impact = likelihood_impact(f)
            lines.append(
                f"| {i} | {f.get('severity', 'info').upper()} | "
                f"{f.get('confidence', 0):.0%} | {f.get('title', '')} | "
                f"{mitre} | {likelihood} | {impact} |"
            )

        path = self.output_dir / "attack_surface.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)

    # ── Findings JSON ──────────────────────────────────────────────────────────

    def _findings_json(self, state: dict[str, Any]) -> str:
        """Raw findings JSON with full provenance."""
        path = self.output_dir / "findings.json"
        findings = {
            "campaign_id": self.campaign_id,
            "scope_hash": self.scope_hash,
            "nexusrecon_version": self.nexusrecon_version,
            "generated": datetime.utcnow().isoformat(),
            "findings": state.get("findings", []),
        }
        path.write_text(json.dumps(findings, indent=2, default=str), encoding="utf-8")
        return str(path)

    # ── Campaign Metadata ──────────────────────────────────────────────────────

    def _campaign_meta(self, state: dict[str, Any]) -> str:
        """Campaign metadata JSON — scope hash, timestamps, stats."""
        path = self.output_dir / "campaign_meta.json"
        meta = {
            "campaign_id": self.campaign_id,
            "engagement_id": self.engagement_id,
            "scope_hash": self.scope_hash,
            "nexusrecon_version": self.nexusrecon_version,
            "generated": datetime.utcnow().isoformat(),
            "phases_completed": state.get("completed_phases", []),
            "total_findings": len(state.get("findings", [])),
            "total_subdomains": len(state.get("subdomain_intel", {})),
            "total_emails": len(state.get("email_intel", {}).get("emails", {})),
            "report_paths": self.report_paths,
        }
        path.write_text(json.dumps(meta, indent=2, default=str), encoding="utf-8")
        return str(path)

    # ── People / Identity Map ──────────────────────────────────────────────────

    def _people_map(self, state: dict[str, Any]) -> str:
        """Identity-focused report: people, roles, org structure, social profiles."""
        emails = state.get("email_intel", {}).get("emails", {})
        agent_messages = [m for m in state.get("agent_messages", []) if m.get("phase") in ("phase2", "phase4")]

        lines = [
            "# People & Identity Map",
            "",
            f"**Campaign:** {self.campaign_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            "",
            "## Overview",
            "",
            f"Total identities discovered: {len(emails)}",
            "",
        ]

        # Group by department/role
        by_dept: dict[str, list] = {}
        by_role: dict[str, list] = {}
        for em, info in emails.items():
            if isinstance(info, dict):
                dept = info.get("department", "unknown")
                pos = info.get("position", "unknown")
                by_dept.setdefault(dept, []).append((em, info))
                by_role.setdefault(pos, []).append((em, info))

        lines.append("## By Department")
        lines.append("")
        for dept, members in sorted(by_dept.items()):
            lines.append(f"### {dept} ({len(members)})")
            lines.append("")
            for em, info in members:
                pos = info.get("position", "")
                lines.append(f"- `{em}` — {pos}")
            lines.append("")

        lines.append("## By Role/Position")
        lines.append("")
        for role, members in sorted(by_role.items(), key=lambda x: len(x[1]), reverse=True):
            if role in ("unknown", "None", ""):
                continue
            lines.append(f"### {role}")
            lines.append("")
            for em, _ in members[:10]:
                lines.append(f"- `{em}`")
            if len(members) > 10:
                lines.append(f"- ... and {len(members) - 10} more")
            lines.append("")

        # Executive targets
        exec_keywords = ["ceo", "cfo", "cto", "ciso", "cso", "vp", "director", "executive", "president", "founder"]
        executives = []
        for em, info in emails.items():
            if isinstance(info, dict):
                pos = str(info.get("position", "")).lower()
                dept = str(info.get("department", "")).lower()
                if any(kw in pos or kw in dept for kw in exec_keywords):
                    executives.append((em, info))

        if executives:
            lines.append("## Executive Targets (High-Value)")
            lines.append("")
            for em, info in executives:
                lines.append(f"- **{em}** — {info.get('position', 'N/A')} | {info.get('department', 'N/A')}")
            lines.append("")

        # Agent analysis
        if agent_messages:
            lines.append("## Analyst Notes")
            lines.append("")
            for msg in agent_messages:
                lines.append(f"### {msg.get('agent', 'unknown')} — {msg.get('phase', '')}")
                lines.append("")
                lines.append(strip_agent_scaffolding(msg.get("analysis", "")))
                lines.append("")

        path = self.output_dir / "people_identity_map.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)

    # ── Vulnerability Correlation ──────────────────────────────────────────────

    def _vuln_correlation(self, state: dict[str, Any]) -> str:
        """CVE-to-asset mapping, KEV correlation, exploit availability."""
        vuln_intel = state.get("vuln_intel", {})
        findings = [f for f in state.get("findings", []) if f.get("category", "").startswith("vuln")]

        lines = [
            "# Vulnerability Correlation Report",
            "",
            f"**Campaign:** {self.campaign_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            "",
            "## Summary",
            "",
            f"Vulnerability-related findings: {len(findings)}",
            f"Vulnerability intelligence sources: {len(vuln_intel)}",
            "",
        ]

        # KEV findings
        kev_data = vuln_intel.get("kev", {})
        if kev_data:
            lines.append("## Known Exploited Vulnerabilities (CISA KEV)")
            lines.append("")
            kev_entries = kev_data.get("data", {}).get("entries", kev_data.get("vulnerabilities", []))
            if isinstance(kev_entries, list):
                for entry in kev_entries[:20]:
                    if isinstance(entry, dict):
                        lines.append(f"- **{entry.get('cveID', entry.get('cve', 'N/A'))}**: {entry.get('shortDescription', entry.get('description', ''))[:200]}")
                        lines.append(f"  - Vendor: {entry.get('vendorProject', entry.get('vendor', 'N/A'))}")
                        lines.append(f"  - Product: {entry.get('product', 'N/A')}")
                        if entry.get("dateAdded"):
                            lines.append(f"  - KEV Added: {entry['dateAdded']}")
                        lines.append("")

        # NVD findings
        nvd_entries = {k: v for k, v in vuln_intel.items() if k.startswith("nvd/")}
        if nvd_entries:
            lines.append("## NVD CVE Matches by Technology")
            lines.append("")
            for key, data in nvd_entries.items():
                tech = key.replace("nvd/", "")
                lines.append(f"### {tech}")
                lines.append("")
                if isinstance(data, dict):
                    cves = data.get("data", {}).get("vulnerabilities", data.get("cves", []))
                    if isinstance(cves, list):
                        for cve in cves[:10]:
                            if isinstance(cve, dict):
                                cvss = cve.get("cvss_score", cve.get("cvssV3", {}).get("baseScore", "N/A"))
                                desc = str(cve.get("description", "") or "")[:150]
                                lines.append(f"- [{cve.get('id', cve.get('cveId', 'N/A'))}] (CVSS: {cvss}) — {desc}")
                        lines.append("")

        # Correlated findings
        if findings:
            lines.append("## Correlated Findings")
            lines.append("")
            for i, f in enumerate(findings, 1):
                lines.append(f"### {i}. {f.get('title', 'Untitled')}")
                lines.append(f"- Severity: {f.get('severity', 'info').upper()}")
                lines.append(f"- Confidence: {f.get('confidence', 0):.0%}")
                lines.append(f"- Assets: {', '.join(f.get('affected_assets', [])[:5])}")
                if f.get("mitre_techniques"):
                    lines.append(f"- MITRE: {', '.join(f['mitre_techniques'])}")
                lines.append("")

        path = self.output_dir / "vulnerability_correlation.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)

    # ── Vendor / Supply Chain ──────────────────────────────────────────────────

    def _vendor_supply_chain(self, state: dict[str, Any]) -> str:
        """Third-party risk, supply chain exposure, vendor analysis."""
        cloud_intel = state.get("cloud_intel", {})
        code_intel = state.get("code_intel", {})
        subdomain_intel = state.get("subdomain_intel", {})
        infra_intel = state.get("infra_intel", {})

        lines = [
            "# Vendor & Supply Chain Report",
            "",
            f"**Campaign:** {self.campaign_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            "",
            "## Third-Party Services Detected",
            "",
        ]

        # Identify cloud providers from POSITIVE evidence only (F-B4). The
        # recon tools write a cloud_intel entry even when they find nothing,
        # so keying on entry existence would falsely "detect" AWS/GCP/M365.
        providers = set()
        for key, data in cloud_intel.items():
            if not _provider_has_evidence(data):
                continue
            kl = key.lower()
            if "aws" in kl:
                providers.add("AWS")
            if "azure" in kl or "m365" in kl:
                providers.add("Microsoft 365 / Azure")
            if "gcp" in kl:
                providers.add("Google Cloud")

        for provider in sorted(providers):
            lines.append(f"- **{provider}**")
        if not providers:
            lines.append("- No third-party cloud providers confirmed from collected data")
        lines.append("")

        # CDN and hosting detection
        cdn_indicators = ["cloudflare", "fastly", "akamai", "cloudfront", "azurefd", "gcdn"]
        hosting_indicators = ["aws", "azure", "google", "digitalocean", "linode", "hetzner"]

        cdn_found = set()
        hosting_found = set()
        for sub, data in subdomain_intel.items():
            if isinstance(data, dict):
                sources = data.get("sources", [])
                for s in sources:
                    s_lower = str(s).lower()
                    for cdn in cdn_indicators:
                        if cdn in s_lower:
                            cdn_found.add(cdn)
                    for host in hosting_indicators:
                        if host in s_lower:
                            hosting_found.add(host)

        lines.append("## CDN Providers")
        lines.append("")
        for cdn in sorted(cdn_found) if cdn_found else ["None detected"]:
            lines.append(f"- {cdn}")
        lines.append("")

        lines.append("## Hosting Providers")
        lines.append("")
        for host in sorted(hosting_found) if hosting_found else ["None detected"]:
            lines.append(f"- {host}")
        lines.append("")

        # GitHub/Code supply chain. List a source only when it actually
        # surfaced something (F-B4): the run's github_recon / gitleaks /
        # trufflehog all returned zero, yet were listed as "code sources"
        # purely because they were queried.
        code_sources = [k for k, v in code_intel.items() if _code_source_has_evidence(v)]
        if code_sources:
            lines.append("## Code & Package Sources")
            lines.append("")
            for src in code_sources:
                lines.append(f"- {src}")
            lines.append("")

            # Check for leaked dependencies
            for key, data in code_intel.items():
                if isinstance(data, dict):
                    deps = data.get("data", {}).get("dependencies", data.get("packages", []))
                    if deps:
                        lines.append(f"### Dependencies in {key}")
                        lines.append("")
                        for dep in deps[:20]:
                            if isinstance(dep, dict):
                                lines.append(f"- {dep.get('name', 'unknown')} {dep.get('version', '')}")
                            else:
                                lines.append(f"- {dep}")
                        lines.append("")

        # Infrastructure vendors
        lines.append("## Infrastructure Vendors")
        lines.append("")
        for key, data in infra_intel.items():
            if isinstance(data, dict):
                tech = data.get("data", {}).get("tech", data.get("technologies", []))
                if tech:
                    lines.append(f"### {key}")
                    lines.append("")
                    for t in tech[:10]:
                        if isinstance(t, dict):
                            lines.append(f"- {t.get('name', 'unknown')} ({t.get('version', 'unknown')})")
                        else:
                            lines.append(f"- {t}")
                    lines.append("")

        path = self.output_dir / "vendor_supply_chain.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)

    # ── Jira Tracker CSV ───────────────────────────────────────────────────────

    def _jira_tracker(self, state: dict[str, Any]) -> str:
        """Export findings as Jira-compatible CSV for issue tracking."""
        import csv
        path = self.output_dir / "jira_tracker.csv"
        findings = state.get("findings", [])

        severity_to_priority = {
            "critical": "Highest",
            "high": "High",
            "medium": "Medium",
            "low": "Low",
            "info": "Lowest",
        }

        with open(path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow([
                "Summary", "Description", "Priority", "Labels",
                "Components", "Environment", "Affected Assets",
                "MITRE Technique", "Source", "Confidence",
                "Recommendation", "Evidence",
            ])
            for finding in findings:
                writer.writerow([
                    finding.get("title", "Untitled Finding"),
                    finding.get("description", ""),
                    severity_to_priority.get(finding.get("severity", "info"), "Lowest"),
                    f"nexusrecon,{finding.get('category', 'general')},{finding.get('severity', 'info')}",
                    finding.get("category", "general"),
                    self.campaign_id,
                    ", ".join(finding.get("affected_assets", [])[:5]),
                    ", ".join(finding.get("mitre_techniques", [])),
                    finding.get("source", "unknown"),
                    f"{finding.get('confidence', 0):.0%}",
                    finding.get("recommendation", ""),
                    finding.get("raw_evidence_hash", ""),
                ])

        return str(path)

    # ── Entity Graph HTML (pyvis) ──────────────────────────────────────────────

    def _entity_graph_html(self, state: dict[str, Any]) -> str:
        """Interactive HTML entity graph using pyvis."""
        from networkx import DiGraph
        from pyvis.network import Network

        entity_graph = state.get("entity_graph", {})
        if not entity_graph:
            # Create a minimal graph from available data
            entity_graph = {
                "subdomains": list(state.get("subdomain_intel", {}).keys())[:100],
                "emails": list(state.get("email_intel", {}).get("emails", {}).keys())[:100],
            }

        G = DiGraph()

        # Add root target
        seeds = state.get("seeds", [])
        for seed in seeds:
            G.add_node(seed, label=seed, color="#ff0000", size=30, group="seed")

        # Add subdomains
        for sub in entity_graph.get("subdomains", [])[:200]:
            G.add_node(sub, label=sub, color="#0066cc", size=10, group="subdomain")
            for seed in seeds:
                if seed in sub:
                    G.add_edge(seed, sub)

        # Add emails
        for em in entity_graph.get("emails", [])[:200]:
            G.add_node(em, label=em, color="#009900", size=8, group="email")
            domain = em.split("@")[-1] if "@" in em else ""
            for seed in seeds:
                if domain and seed in domain:
                    G.add_edge(seed, em)

        # Add confirmed leads as nodes
        for lead in state.get("confirmed_leads", [])[:20]:
            lead_node = f"lead:{lead[:50]}"
            G.add_node(lead_node, label=lead[:50], color="#ff6600", size=12, group="lead")
            for seed in seeds:
                G.add_edge(seed, lead_node)

        net = Network(height="750px", width="100%", directed=True, notebook=False)
        net.from_nx(G)

        # Style
        net.set_options("""
        {
          "physics": {
            "enabled": true,
            "solver": "forceAtlas2Based",
            "forceAtlas2Based": {"gravitationalConstant": -50, "centralGravity": 0.01, "springLength": 100}
          },
          "interaction": {"hover": true, "tooltipDelay": 200}
        }
        """)

        path = self.output_dir / "entity_graph.html"
        net.write_html(str(path))
        return str(path)

    # ── PDF Report ─────────────────────────────────────────────────────────────

    def _pdf_report(self, state: dict[str, Any]) -> str:
        """Generate PDF executive summary report."""
        try:
            from weasyprint import HTML
        except ImportError:
            log.warning("weasyprint not installed — generating PDF as HTML fallback")
            try:
                from rich.console import Console as _Console
                _Console().print(
                    "[yellow]PDF generation requires weasyprint: "
                    "pip install nexusrecon[pdf][/yellow]"
                )
            except Exception:
                pass
            return self._pdf_report_fallback(state)

        findings = state.get("findings", [])
        critical = [f for f in findings if f.get("severity") == "critical"]
        high = [f for f in findings if f.get("severity") == "high"]

        def _severity_css(s: str) -> str:
            return s.lower() if s in ("critical", "high", "medium") else "info"

        top_rows = ""
        for i, f in enumerate(findings[:15], 1):
            sev = f.get("severity", "info").upper()
            sev_cls = _severity_css(f.get("severity", "info"))
            title = (f.get("title", "") or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            conf = f"{f.get('confidence', 0):.0%}"
            top_rows += f"<tr><td>{i}</td><td class=\"{sev_cls}\">{sev}</td><td>{title}</td><td>{conf}</td></tr>"

        leads_html = ""
        for lead in state.get("confirmed_leads", [])[:20]:
            safe = str(lead).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            leads_html += f"<li>{safe}</li>"

        html_content = f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
                h1 {{ color: #1a1a2e; border-bottom: 3px solid #e94560; padding-bottom: 10px; }}
                h2 {{ color: #16213e; margin-top: 30px; }}
                .critical {{ color: #e94560; font-weight: bold; }}
                .high {{ color: #ff6b35; font-weight: bold; }}
                .medium {{ color: #f7b731; }}
                .meta {{ color: #666; font-size: 0.9em; }}
                table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
                th, td {{ border: 1px solid #ddd; padding: 10px; text-align: left; }}
                th {{ background-color: #1a1a2e; color: white; }}
                tr:nth-child(even) {{ background-color: #f8f9fa; }}
                .footer {{ margin-top: 40px; font-size: 0.8em; color: #999; text-align: center; }}
            </style>
        </head>
        <body>
            <h1>NexusRecon — Executive Report</h1>
            <p class="meta">
                Campaign: {self.campaign_id}<br>
                Engagement: {self.engagement_id}<br>
                Generated: {datetime.utcnow().isoformat()}<br>
                Scope Hash: {self.scope_hash}
            </p>

            <h2>Overview</h2>
            <table>
                <tr><th>Metric</th><th>Count</th></tr>
                <tr><td>Total Findings</td><td>{len(findings)}</td></tr>
                <tr><td>Critical</td><td class="critical">{len(critical)}</td></tr>
                <tr><td>High</td><td class="high">{len(high)}</td></tr>
                <tr><td>Subdomains</td><td>{len(state.get('subdomain_intel', {}))}</td></tr>
                <tr><td>Emails</td><td>{len(state.get('email_intel', {}).get('emails', {}))}</td></tr>
            </table>

            <h2>Top Findings</h2>
            <table>
                <tr><th>#</th><th>Severity</th><th>Title</th><th>Confidence</th></tr>
                {top_rows}
            </table>

            <h2>Confirmed Attack Leads</h2>
            <ul>
                {leads_html}
            </ul>

            <div class="footer">
                Generated by NexusRecon v1.0.0 — Authorized use only
            </div>
        </body>
        </html>
        """

        path = self.output_dir / "report.pdf"
        HTML(string=html_content).write_pdf(str(path))
        return str(path)

    def _pdf_report_fallback(self, state: dict[str, Any]) -> str:
        """Fallback: generate HTML that can be printed to PDF manually."""
        findings = state.get("findings", [])
        lines = [
            "<!DOCTYPE html><html><head><title>NexusRecon Report</title>",
            "<style>body{font-family:monospace;margin:40px;}</style>",
            "</head><body>",
            f"<h1>NexusRecon Report — {self.campaign_id}</h1>",
            f"<p>Generated: {datetime.utcnow().isoformat()}</p>",
            f"<p>Total findings: {len(findings)}</p>",
            "<h2>Findings</h2><ul>",
        ]
        for f in findings[:30]:
            lines.append(f'<li><b>[{f.get("severity", "info").upper()}]</b> {f.get("title", "")}</li>')
        lines.append("</ul></body></html>")

        path = self.output_dir / "report.html"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)

    # ── PPTX Report ────────────────────────────────────────────────────────────

    def _pptx_report(self, state: dict[str, Any]) -> str:
        """Generate PowerPoint executive briefing."""
        try:
            from pptx import Presentation
            from pptx.util import Emu, Pt
        except ImportError:
            log.warning("python-pptx not installed — skipping PPTX generation")
            return ""

        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)

        findings = state.get("findings", [])
        critical = [f for f in findings if f.get("severity") == "critical"]
        high = [f for f in findings if f.get("severity") == "high"]

        def add_slide(title: str, bullets: list[str]):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(14)
                p.level = 0

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "NexusRecon — Executive Briefing"
        title_slide.placeholders[1].text = (
            f"Campaign: {self.campaign_id}\n"
            f"Engagement: {self.engagement_id}\n"
            f"Generated: {datetime.utcnow().strftime('%Y-%m-%d')}"
        )

        # Overview
        add_slide("Overview", [
            f"Total findings: {len(findings)}",
            f"Critical severity: {len(critical)}",
            f"High severity: {len(high)}",
            f"Subdomains discovered: {len(state.get('subdomain_intel', {}))}",
            f"Emails discovered: {len(state.get('email_intel', {}).get('emails', {}))}",
            f"Phases completed: {', '.join(state.get('completed_phases', []))}",
        ])

        # Top findings
        top_findings_bullets = []
        for i, f in enumerate(findings[:8], 1):
            top_findings_bullets.append(
                f"{i}. [{f.get('severity', 'info').upper()}] {f.get('title', '')}"
            )
        add_slide("Top Findings", top_findings_bullets if top_findings_bullets else ["No findings recorded"])

        # Attack leads
        leads = state.get("confirmed_leads", [])
        add_slide("Confirmed Attack Leads", leads[:10] if leads else ["No confirmed leads"])

        # Recommendations
        recommendations = []
        if critical:
            recommendations.append("IMMEDIATE: Address critical findings within 24 hours")
        if high:
            recommendations.append("URGENT: Remediate high-severity findings within 7 days")
        recommendations.append("Implement continuous monitoring for exposed assets")
        recommendations.append("Review and rotate any exposed credentials")
        recommendations.append("Update email security policies (DMARC/SPF/DKIM)")
        add_slide("Recommendations", recommendations)

        path = self.output_dir / "executive_briefing.pptx"
        prs.save(str(path))
        return str(path)

    # ── Harvested Credentials Report ──────────────────────────────────────────

    def _harvested_credentials(self, state: dict[str, Any]) -> str:
        creds = state.get("harvested_credentials", [])

        # F-B8: no false-alarm "contains real credentials" banner on an empty
        # file. The header only fires when there is actually something secret
        # to protect; an empty run says so plainly.
        if not creds:
            lines = [
                "# Harvested Credentials",
                "",
                f"**Campaign:** {self.campaign_id}",
                f"**Generated:** {datetime.utcnow().isoformat()}",
                "",
                "No credentials were harvested. File retained for completeness.",
                "",
            ]
            md_path = self.output_dir / "harvested_credentials.md"
            md_path.write_text("\n".join(lines), encoding="utf-8")
            json_path = self.output_dir / "harvested_credentials.json"
            json_path.write_text(
                json.dumps({"campaign_id": self.campaign_id, "credentials": []}, indent=2, default=str),
                encoding="utf-8",
            )
            return str(md_path)

        lines = [
            "⚠ This file contains real credentials. Treat as Secret. Rotate before sharing.",
            "",
            "# Harvested Credentials",
            "",
            f"**Campaign:** {self.campaign_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            f"**Total:** {len(creds)} credential(s)",
            "",
        ]

        from collections import defaultdict
        grouped: dict[str, dict[str, list]] = defaultdict(lambda: defaultdict(list))
        for cred in creds:
            grouped[cred.get("source_type", "unknown")][cred.get("cred_type", "unknown")].append(cred)

        for source_type, by_type in sorted(grouped.items()):
            lines.append(f"## Source: {source_type}")
            lines.append("")
            for cred_type, items in sorted(by_type.items()):
                lines.append(f"### {cred_type} ({len(items)})")
                lines.append("")
                for c in items:
                    lines.append(f"- **Redacted:** `{c.get('value_redacted', '?')}`")
                    lines.append(f"  - Source: {c.get('source_url', '?')}")
                    lines.append(f"  - Validated: {'✓' if c.get('validated') else '✗'}")
                    if c.get("validation_metadata"):
                        lines.append(f"  - Metadata: `{str(c['validation_metadata'])[:200]}`")
                    if c.get("next_steps"):
                        for step in c["next_steps"]:
                            lines.append(f"  - Next: {step}")
                    lines.append("")

        md_path = self.output_dir / "harvested_credentials.md"
        md_path.write_text("\n".join(lines), encoding="utf-8")

        json_path = self.output_dir / "harvested_credentials.json"
        json_path.write_text(
            json.dumps({"campaign_id": self.campaign_id, "credentials": creds}, indent=2, default=str),
            encoding="utf-8",
        )
        return str(md_path)

    # ── Credential Exposure Paths (D7) ────────────────────────────────────────

    def _credential_exposure_paths(self, state: dict[str, Any]) -> str:
        """Operator punch list: personal breach data → corporate auth surfaces.

        This document is the primary deliverable from Phase D (identity
        attribution + credential correlation). It bridges personal-identity
        breach data to testable corporate authentication endpoints.

        OPERATOR NOTICE — FOR YOUR EYES ONLY
        Every candidate requires explicit operator review and authorisation.
        This framework never executes any test automatically.
        """

        punch_list: list[dict[str, Any]] = state.get("credential_punch_list", [])
        graph_data: dict[str, Any] = state.get("identity_graph", {})
        pivot_results: dict[str, Any] = state.get("personal_pivot_results", {})

        # ── Header + notice ───────────────────────────────────────────────
        lines: list[str] = [
            "# Credential Exposure Paths",
            "",
            f"**Campaign:** {self.campaign_id}",
            f"**Engagement:** {self.engagement_id}",
            f"**Generated:** {datetime.utcnow().isoformat()}",
            "",
            "> **OPERATOR NOTICE — CONFIDENTIAL**",
            ">",
            "> This document maps personal credential exposures to corporate",
            "> authentication surfaces. Every entry is a CANDIDATE for human",
            "> review — not an instruction to execute.",
            ">",
            "> **DO NOT automate credential testing. DO NOT share outside the",
            "> engagement team. Rotate or destroy after the engagement.**",
            "",
            "---",
            "",
        ]

        # ── Executive summary ─────────────────────────────────────────────
        identity_count = graph_data.get("identity_count", 0)
        pivoted_count = len(pivot_results)

        if punch_list:
            # Re-use the summarise utility without importing live modules
            # (the list is already serialised dicts at this point).
            by_kind: dict[str, int] = {}
            by_ep: dict[str, int] = {}
            by_band = {"high": 0, "medium": 0, "low": 0}
            mfa_count = 0
            identity_ids: set = set()

            for c in punch_list:
                k = c.get("credential_kind", "unknown")
                by_kind[k] = by_kind.get(k, 0) + 1
                ep = c.get("endpoint_type", "unknown")
                by_ep[ep] = by_ep.get(ep, 0) + 1
                identity_ids.add(c.get("identity_id", ""))
                if c.get("mfa_expected"):
                    mfa_count += 1
                conf = c.get("confidence", 0.0)
                if conf >= 0.70:
                    by_band["high"] += 1
                elif conf >= 0.40:
                    by_band["medium"] += 1
                else:
                    by_band["low"] += 1

            lines += [
                "## Executive Summary",
                "",
                f"Phase 2.5 pivoted **{pivoted_count}** of **{identity_count}** "
                f"discovered identities to their personal identity. The correlation "
                f"engine produced **{len(punch_list)}** candidate credential-test "
                f"entries across **{len(identity_ids)}** unique identities.",
                "",
                "### Candidate summary",
                "",
                "| Metric | Value |",
                "|--------|-------|",
                f"| Total candidates | {len(punch_list)} |",
                f"| Identities with exposure | {len(identity_ids)} |",
                f"| High-confidence (≥70%) | {by_band['high']} |",
                f"| Medium-confidence (40–69%) | {by_band['medium']} |",
                f"| Low-confidence (<40%) | {by_band['low']} |",
                f"| MFA-protected endpoints | {mfa_count} |",
                "",
                "**By credential kind:**",
                "",
            ]
            for kind, cnt in sorted(by_kind.items(), key=lambda x: -x[1]):
                lines.append(f"- `{kind}`: {cnt}")
            lines.append("")
            lines.append("**By auth endpoint type:**")
            lines.append("")
            for ep, cnt in sorted(by_ep.items(), key=lambda x: -x[1]):
                lines.append(f"- `{ep}`: {cnt}")
            lines += ["", "---", ""]

        else:
            lines += [
                "## Executive Summary",
                "",
                "Phase 2.5 found **no actionable credential exposure paths**.",
                "",
                f"Identities pivoted: **{pivoted_count}** / **{identity_count}**",
                "",
                "Possible causes:",
                "- No personal email candidates confirmed by breach DBs",
                "- No auth endpoints discovered in cloud_intel (Phase 2 needed)",
                "- Breach DB tools not available (check HIBP, DeHashed, "
                "HudsonRock API keys)",
                "",
                "---",
                "",
            ]

        # ── Punch list ────────────────────────────────────────────────────

        if punch_list:
            lines += [
                "## Punch List",
                "",
                "> Sorted by confidence descending. All `credential_value` fields",
                "> are redacted. Do not execute any test without explicit",
                "> operator authorisation and confirmed scope coverage.",
                "",
            ]

            _band_emoji = {True: "🔴", False: ""}

            for rank, candidate in enumerate(punch_list, 1):
                conf = candidate.get("confidence", 0.0)
                conf_band = (
                    "HIGH" if conf >= 0.70 else
                    "MEDIUM" if conf >= 0.40 else "LOW"
                )
                mfa_note = " ⚠ MFA" if candidate.get("mfa_expected") else ""
                identity_label = candidate.get("identity_label", "Unknown")
                corp_email = candidate.get("corp_email", "")
                breach_src = candidate.get("breach_source", "")
                cred_kind = candidate.get("credential_kind", "")
                breach_date = candidate.get("breach_date") or "date unknown"
                observed_at = candidate.get("observed_at", "")
                ep_url = candidate.get("test_endpoint_url", "")
                ep_type = candidate.get("endpoint_type", "")
                risk_flags = candidate.get("risk_flags", [])
                mitre = candidate.get("mitre_techniques", [])
                notes = candidate.get("notes", "")

                lines += [
                    f"### [{rank}] [{conf_band}]{mfa_note} {identity_label}",
                    "",
                    "| Field | Value |",
                    "|-------|-------|",
                    f"| **Corp email** | `{corp_email}` |",
                    f"| **Breach** | {breach_src} ({breach_date}) |",
                    f"| **Observed at** | `{observed_at}` |",
                    f"| **Credential kind** | `{cred_kind}` |",
                    f"| **Breach confidence** | {candidate.get('breach_confidence', '')} |",
                    "| **Credential value** | `[REDACTED]` |",
                    f"| **Test endpoint** | `{ep_url}` |",
                    f"| **Endpoint type** | `{ep_type}` |",
                    f"| **Confidence** | {conf:.0%} |",
                    "",
                ]

                if risk_flags:
                    lines.append("**Risk flags:**")
                    lines.append("")
                    for flag in risk_flags:
                        lines.append(f"- `{flag}`")
                    lines.append("")

                if mitre:
                    mitre_str = ", ".join(
                        f"[{t}](https://attack.mitre.org/techniques/{t.replace('.', '/')})"
                        for t in mitre
                    )
                    lines += [f"**MITRE:** {mitre_str}", ""]

                if notes:
                    lines += [f"**Notes:** {notes}", ""]

                lines += [
                    "> ⛔ OPERATOR DECISION REQUIRED — do not execute without",
                    "> explicit authorisation and confirmed scope.",
                    "",
                    "---",
                    "",
                ]

        # ── Identity graph summary ────────────────────────────────────────

        identities_data = graph_data.get("identities", [])
        if identities_data:
            lines += [
                "## Identity Graph Summary",
                "",
                f"Total identities: **{identity_count}**",
                "",
                "| Identity | Corp email | Personal emails | Handles | "
                "Credential exposures |",
                "|----------|-----------|-----------------|---------|---------------------|",
            ]
            for ident in identities_data[:30]:
                label = ident.get("primary_label", "Unknown")[:40]
                idents = ident.get("identifiers", [])
                corp_emails = [
                    i["value"] for i in idents
                    if i.get("identifier_type") == "corp_email"
                ]
                personal_emails = [
                    i["value"] for i in idents
                    if i.get("identifier_type") == "personal_email"
                ]
                handles = [
                    i["value"] for i in idents
                    if i.get("identifier_type") == "handle"
                ]
                exposures = len(ident.get("credential_exposures", []))
                corp_str = corp_emails[0] if corp_emails else "—"
                lines.append(
                    f"| {label} | `{corp_str}` | {len(personal_emails)} | "
                    f"{len(handles)} | {exposures} |"
                )
            lines += ["", "---", ""]

        # ── Footer ────────────────────────────────────────────────────────

        lines += [
            "## Legal + Operational Disclaimer",
            "",
            "This document was generated by NexusRecon as part of an "
            "authorised penetration testing engagement. All credential "
            "information was obtained from publicly available breach "
            "databases and infostealer intelligence sources.",
            "",
            "- All testing must remain within the confirmed engagement scope.",
            "- Credentials must not be used for unauthorised access.",
            "- This document must be destroyed after engagement close-out.",
            "- Never automate credential submission without explicit written "
            "authorisation.",
            "",
        ]

        path = self.output_dir / "credential_exposure_paths.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        log.info(
            "Credential exposure paths report written",
            path=str(path),
            candidate_count=len(punch_list),
        )

        # Also emit the punch list as JSON for downstream tool consumption.
        json_path = self.output_dir / "credential_punch_list.json"
        json_path.write_text(
            json.dumps(
                {
                    "campaign_id": self.campaign_id,
                    "engagement_id": self.engagement_id,
                    "generated": datetime.utcnow().isoformat(),
                    "do_not_execute": True,
                    "candidates": punch_list,
                },
                indent=2,
                default=str,
            ),
            encoding="utf-8",
        )

        return str(path)

    # ── Spear-Phishing Intelligence (Phase E11) ────────────────────────────────

    def _master_report_obsidian(self, state: dict[str, Any]) -> str:
        """Emit ``master_report.obsidian.md`` next to the standard
        master report.

        Reads the just-written ``master_report.md`` from disk and
        runs it through the Obsidian transforms (frontmatter,
        wikilinks, callouts). This way the prose generation has
        exactly one source — a future edit to the master_report
        builder picks up automatically in the Obsidian variant.
        """
        from nexusrecon.reports.obsidian_export import build_obsidian_master

        standard_path = Path(self.report_paths["master_report"])
        standard_md = standard_path.read_text(encoding="utf-8")
        out = build_obsidian_master(
            standard_md=standard_md,
            state={
                **state,
                "campaign_id": self.campaign_id,
                "engagement_id": self.engagement_id,
                "generated": datetime.utcnow().isoformat(),
            },
            scope_hash=self.scope_hash,
            nexusrecon_version=self.nexusrecon_version,
        )
        path = self.output_dir / "master_report.obsidian.md"
        path.write_text(out, encoding="utf-8")
        return str(path)

    def _spear_phishing_intelligence(
        self, state: dict[str, Any],
    ) -> tuple[str, str]:
        """Phase E11 deliverable: per-target pretext dossiers + raw
        candidate JSON. Always emits both files (empty content when
        no candidates surfaced) so the path is stable for the master
        report's links."""
        from nexusrecon.reports.spear_phishing_intelligence import (
            build_spear_phishing_intelligence_md,
        )

        return build_spear_phishing_intelligence_md(
            campaign_id=self.campaign_id,
            engagement_id=self.engagement_id,
            state=state,
            output_dir=self.output_dir,
        )

    # ── Master Report (V3 Move 2) ──────────────────────────────────────────────

    def _render_run_health_block(self, state: dict[str, Any]) -> list[str]:
        """Master-report run-health banner (Wave F-A5).

        Surfaces whether the run can be trusted *before* the findings are
        read: the blunt caveats first, then a one-line tool-outcome tally,
        the degraded capabilities, and the analysis-engine provenance
        (live model vs. MockLLM fallback). Reads ``state["run_health"]``,
        which run_campaign populates before report generation.
        """
        rh = state.get("run_health") or {}
        if not rh:
            return []
        out = ["## 1a. Run Health", ""]
        for c in rh.get("caveats", []):
            out.append(f"> {c}")
        if rh.get("caveats"):
            out.append("")
        out.append(
            f"- **Tool outcomes:** {rh.get('productive', 0)} returned data, "
            f"{len(rh.get('degraded', []))} degraded, {len(rh.get('errors', []))} errored, "
            f"{len(rh.get('policy_skipped', []))} skipped by policy"
        )
        caps = [c.get("capability") for c in rh.get("degraded_capabilities", [])]
        if caps:
            out.append(f"- **Degraded capabilities:** {', '.join(caps)}")
        out.append(
            f"- **Analysis engine:** {rh.get('llm_mode', 'unknown')} "
            f"({rh.get('llm_calls', 0)} call(s), ${rh.get('llm_cost_usd', 0.0):.4f})"
        )
        if rh.get("node_estimate_note"):
            out.append(f"- **Forecast vs. reality:** {rh['node_estimate_note']}")
        out.append("")
        return out

    def _master_report(self, state: dict[str, Any]) -> str:
        """
        Master report — single cohesive narrative deliverable.

        Layout:
          1. Snapshot (engine, always)
          2. Executive Brief (agent prose)
          3. Top Threads to Pull (engine, embeds top_threads.md)
          4-8. Attack Surface / Personas / Vulns / Creds / Pretext (agent prose,
               only sections with content)
          9. Evidence & Provenance (engine, always)
          10. Recommendations (agent prose)
          11. Appendix — Deeper Reading (engine, dynamic by file existence)
        """
        emails = state.get("email_intel", {}).get("emails", {})
        cloud_intel = state.get("cloud_intel", {})
        code_intel = state.get("code_intel", {})
        subdomain_intel = state.get("subdomain_intel", {})
        infra_intel = state.get("infra_intel", {})
        vuln_intel = state.get("vuln_intel", {})
        pretext_intel = state.get("pretext_intel", {})
        harvested_creds = state.get("harvested_credentials", [])
        findings = state.get("findings", [])
        ranked_threads = state.get("ranked_threads", [])

        def _verified(d: Any) -> bool:
            # V3 Move 2 skip-empty rule: a cloud entry counts as "verified"
            # only when the tool explicitly tagged attribution_confidence >= 0.5.
            # Untagged data (e.g., legacy or stub tool output) defaults to 0.0
            # so we don't render the Cloud subsection on noise alone.
            if not isinstance(d, dict):
                return False
            ac = d.get("attribution_confidence")
            try:
                return ac is not None and float(ac) >= 0.5
            except (TypeError, ValueError):
                return False

        def _code_has_findings(intel: Any) -> bool:
            """True iff any code-tool result has actual leaks/findings/repos."""
            if not isinstance(intel, dict):
                return False
            for v in intel.values():
                if not isinstance(v, dict):
                    continue
                for key in ("leaks", "findings", "dork_results", "secrets", "repos"):
                    if v.get(key):
                        return True
                if isinstance(v.get("org_repos"), dict) and v["org_repos"].get("total", 0) > 0:
                    return True
            return False

        def _infra_has_findings(intel: Any) -> bool:
            """True iff any infra tool surfaced at least one host/result."""
            if not isinstance(intel, dict):
                return False
            for v in intel.values():
                if not isinstance(v, dict):
                    continue
                if isinstance(v.get("search"), dict) and v["search"].get("total", 0) > 0:
                    return True
                for key in ("hosts", "findings", "results"):
                    if v.get(key):
                        return True
            return False

        has_verified_cloud_oidc = any(
            _verified(d) and isinstance(d.get("openid_config"), dict) and d["openid_config"].get("found")
            for d in cloud_intel.values()
        )
        has_identity = len(emails) > 0 or has_verified_cloud_oidc
        has_cloud = any(_verified(d) for d in cloud_intel.values())
        has_code = _code_has_findings(code_intel)
        has_network = len(subdomain_intel) > 0 or _infra_has_findings(infra_intel)
        has_personas = len(emails) > 0
        has_vulns = bool(
            vuln_intel.get("enriched_cves") if isinstance(vuln_intel, dict) else {}
        ) or bool(
            (vuln_intel.get("nuclei_scan", {}) or {}).get("findings", [])
            if isinstance(vuln_intel, dict) else False
        )
        has_creds = len(harvested_creds) > 0 if isinstance(harvested_creds, list) else False
        has_pretext = (
            isinstance(pretext_intel, dict)
            and bool(pretext_intel)
            and any(v for v in pretext_intel.values())
        )
        has_attack_surface = has_identity or has_cloud or has_code or has_network

        # Trimmed context for the agent
        cloud_summary: list[dict[str, Any]] = []
        for key, data in list(cloud_intel.items())[:8]:
            if not isinstance(data, dict):
                continue
            cloud_summary.append({
                "key": key,
                "attribution_confidence": data.get("attribution_confidence", 1.0),
                "tenant_id": (
                    data.get("tenant_id")
                    or (data.get("openid_config", {}) or {}).get("tenant_id")
                ),
                "summary": str(data.get("summary", ""))[:200],
            })

        finding_titles = [
            {"title": f.get("title"), "severity": f.get("severity"), "phase": f.get("phase")}
            for f in findings[:25]
        ]
        top_threads_preview = [
            {
                "title": t.get("title"),
                "severity": t.get("severity"),
                "score": t.get("score"),
                "category": t.get("category"),
            }
            for t in ranked_threads[:10]
        ]

        # Build the section include/skip lists for the agent prompt
        section_lines = ["## 2. Executive Brief"]
        if has_attack_surface:
            section_lines.append("## 4. Attack Surface at a Glance")
            if has_identity:
                section_lines.append("### 4.1 Identity")
            if has_cloud:
                section_lines.append("### 4.2 Cloud")
            if has_code:
                section_lines.append("### 4.3 Code & Secrets")
            if has_network:
                section_lines.append("### 4.4 Network")
        if has_personas:
            section_lines.append("## 5. Identified Personas")
        if has_vulns:
            section_lines.append("## 6. Vulnerability Correlation")
        if has_creds:
            section_lines.append("## 7. Harvested Credentials")
        if has_pretext:
            section_lines.append("## 8. Pretext & HUMINT")
        section_lines.append("## 10. Recommendations")

        skip_list: list[str] = []
        for name, included in [
            ("Attack Surface (section 4)", has_attack_surface),
            ("Personas (section 5)", has_personas),
            ("Vulnerabilities (section 6)", has_vulns),
            ("Credentials (section 7)", has_creds),
            ("Pretext & HUMINT (section 8)", has_pretext),
        ]:
            if not included:
                skip_list.append(name)

        task_prompt = (
            "You are producing the body of the master_report.md for this OSINT "
            "campaign. Output ONLY these sections, in this exact order, using "
            "exactly these heading levels:\n\n"
            + "\n".join(f"  {s}" for s in section_lines)
            + "\n\nDO NOT produce these sections (they are empty for this "
            "campaign): "
            + (", ".join(skip_list) if skip_list else "(none — all sections apply)")
            + ".\n\nRequirements:\n"
            "- Section 2 (Executive Brief): 200-400 words. Paragraph 1: what "
            "was discovered (scope, scale, principal assets). Paragraph 2: "
            "what it means operationally. Then 3-5 bullet 'key risks' citing "
            "specific finding titles from the state.\n"
            "- Section 4 subsections (if included): ~60-120 words each. Cite "
            "specific numbers from the state above. Any cloud asset with "
            "attribution_confidence < 0.5 MUST be tagged [POSSIBLE] in prose.\n"
            "- Sections 5-8 (if included): 80-150 words. Specific, cite real "
            "values from the state.\n"
            "- Section 10 (Recommendations): 5-10 numbered prioritized actions. "
            "Each one cites a specific finding title or asset from the state. "
            "No generic boilerplate.\n"
            "- Do NOT produce sections 1, 3, 9, or 11 — the engine writes those.\n"
            "- Output pure markdown only. No surrounding code fences, no "
            "commentary outside the requested sections."
        )

        body_md = ""
        try:
            import asyncio

            from nexusrecon.core.config import get_config
            from nexusrecon.graph.agent_executor import AgentExecutor

            executor = AgentExecutor(get_config())
            result = asyncio.run(
                executor.run_agent(
                    "master_reporter",
                    task_data={
                        "campaign_id": self.campaign_id,
                        "engagement_id": self.engagement_id,
                        "seeds": state.get("seeds", []),
                        "subdomain_count": len(subdomain_intel),
                        "email_count": len(emails),
                        "cloud_summary": cloud_summary,
                        "code_sources": (
                            list(code_intel.keys())[:10]
                            if isinstance(code_intel, dict) else []
                        ),
                        "vuln_summary": {
                            "enriched_cves": list(
                                (vuln_intel.get("enriched_cves") if isinstance(vuln_intel, dict) else {}) or {}
                            )[:15],
                            "nuclei_findings": len(
                                (vuln_intel.get("nuclei_scan", {}) or {}).get("findings", [])
                                if isinstance(vuln_intel, dict) else []
                            ),
                        },
                        "harvested_creds_count": len(harvested_creds) if isinstance(harvested_creds, list) else 0,
                        "pretext_keys": (
                            list(pretext_intel.keys())[:10]
                            if isinstance(pretext_intel, dict) else []
                        ),
                        "findings_titles": finding_titles,
                        "top_threads": top_threads_preview,
                        "section_plan": {
                            "include": section_lines,
                            "skip": skip_list,
                        },
                    },
                    task_prompt=task_prompt,
                    state=state,
                )
            )
            body_md = str(result.get("output", "")).strip()
        except Exception as exc:
            log.warning("master_reporter agent failed", error=str(exc))
            body_md = (
                "## 2. Executive Brief\n\n"
                f"*Agent synthesis unavailable ({exc}). The campaign produced "
                f"{len(findings)} findings across "
                f"{len(state.get('completed_phases', []))} phases. See the "
                "linked reports below for detail.*\n\n"
                "## 10. Recommendations\n\n"
                "*Recommendations could not be auto-generated. Review "
                "`top_threads.md` for the ranked attack paths.*"
            )

        # Assemble the master report
        lines: list[str] = []

        # Section 1: Snapshot — always
        sev_counts = {
            "critical": sum(1 for f in findings if f.get("severity") == "critical"),
            "high": sum(1 for f in findings if f.get("severity") == "high"),
            "medium": sum(1 for f in findings if f.get("severity") == "medium"),
            "low": sum(1 for f in findings if f.get("severity") == "low"),
            "info": sum(1 for f in findings if f.get("severity") == "info"),
        }
        top_sev = "info"
        for sev in ("critical", "high", "medium", "low", "info"):
            if sev_counts.get(sev, 0) > 0:
                top_sev = sev
                break

        lines.extend([
            "# Master Report",
            "",
            f"**Campaign ID:** `{self.campaign_id}`  ",
            f"**Engagement ID:** `{self.engagement_id}`  ",
            f"**Generated:** {datetime.utcnow().isoformat()}  ",
            f"**Scope Hash:** `{self.scope_hash}`  ",
            f"**Tooling:** NexusRecon v`{self.nexusrecon_version}`",
            "",
            "---",
            "",
            "## 1. Snapshot",
            "",
            f"- **Findings:** {len(findings)} "
            f"(top severity: **{top_sev.upper()}**; "
            f"crit/high/med = {sev_counts['critical']}/{sev_counts['high']}/{sev_counts['medium']})",
            f"- **Ranked threats:** {len(ranked_threads)}",
            f"- **Subdomains discovered:** {len(subdomain_intel)}",
            f"- **Emails harvested:** {len(emails)}",
            f"- **Phases completed:** {len(state.get('completed_phases', []))}",
            f"- **LLM spend:** ${state.get('llm_cost_usd', 0.0):.2f}",
            "",
            f"> Authorized engagement under SOW `{self.scope_hash}`. "
            f"All tool activity is hash-chained in `logs/audit_log.jsonl`.",
            "",
        ])

        # Section 1a: Run Health — trust banner before the findings (F-A5).
        lines.extend(self._render_run_health_block(state))

        # Section 2 (Executive Brief) + 4-8 (conditional) from agent.
        # Split body_md at "## 10. Recommendations" so we can place that
        # section AFTER Section 9 (Evidence) — preserves the spec'd
        # numerical reading order 1→2→3→4-8→9→10→11.
        body_pre_recs = body_md
        body_recs = ""
        if body_md:
            marker_idx = body_md.find("## 10. Recommendations")
            if marker_idx >= 0:
                body_pre_recs = body_md[:marker_idx].rstrip()
                body_recs = body_md[marker_idx:].rstrip()

        if body_pre_recs:
            lines.extend(["", body_pre_recs, ""])

        # Section 3: Top Threads — embed top_threads.md content if present
        lines.extend(["", "## 3. Top Threads to Pull", ""])
        tt_path = self.output_dir / "top_threads.md"
        if tt_path.exists():
            try:
                tt_text = tt_path.read_text(encoding="utf-8")
                # Drop the H1 from top_threads.md so we don't have a duplicate top-level heading
                tt_stripped = []
                skipped_h1 = False
                for ln in tt_text.split("\n"):
                    if not skipped_h1 and ln.startswith("# "):
                        skipped_h1 = True
                        continue
                    tt_stripped.append(ln)
                lines.append("\n".join(tt_stripped).strip())
                lines.append("")
            except Exception:
                lines.append("*See `top_threads.md` for the full ranked attack-path list.*")
                lines.append("")
        elif ranked_threads:
            for i, t in enumerate(ranked_threads[:10], 1):
                lines.append(
                    f"{i}. **[{str(t.get('severity', 'info')).upper()}]** "
                    f"{t.get('title', 'Untitled')} — priority "
                    f"{float(t.get('score', 0.0)) * 100:.0f}%"
                )
            lines.append("")
        else:
            lines.append(
                "*No ranked threats produced — Phase 8 may not have completed.*"
            )
            lines.append("")

        # Section 9: Evidence & Provenance — always
        tool_errors = state.get("errors", [])
        completed = state.get("completed_phases", [])
        lines.extend([
            "",
            "---",
            "",
            "## 9. Evidence & Provenance",
            "",
            f"- **Scope hash:** `{self.scope_hash}`",
            f"- **Phases completed:** {', '.join(completed) if completed else '(none)'}",
            f"- **Errors during run:** {len(tool_errors)}",
            f"- **LLM cost:** ${state.get('llm_cost_usd', 0.0):.4f}",
            f"- **Tool cost:** ${state.get('tool_cost_usd', 0.0):.4f}",
            "- **Audit chain:** see `logs/audit_log.jsonl` (hash-chained per tool call)",
            "",
        ])

        # Section 10: Recommendations (split out from agent's body_md so it
        # lands in numerical order, after Section 9 Evidence).
        if body_recs:
            lines.extend(["", body_recs, ""])

        # Section 11: Appendix — Deeper Reading (dynamic by file existence)
        lines.extend([
            "---",
            "",
            "## 11. Appendix: Deeper Reading",
            "",
        ])
        appendix_entries = [
            ("findings.json", "Complete findings JSON with provenance hashes"),
            ("top_threads.md", "Top 10 attack paths (full detail)"),
            ("attack_surface.md", "Severity × confidence × MITRE matrix"),
            ("asset_inventory.md", "Complete asset listing"),
            ("cloud_posture.md", "Cloud and federation analysis"),
            ("vulnerability_correlation.md", "CVE-to-asset mapping"),
            ("people_identity_map.md", "Org chart synthesis"),
            ("vendor_supply_chain.md", "Third-party services detected"),
            ("phishing_package.md", "Per-target phishing draft index"),
            ("harvested_credentials.md", "Exposed credentials (redacted)"),
            ("entity_graph.html", "Interactive entity graph"),
            ("jira_tracker.csv", "Findings in Jira import format"),
            ("executive_briefing.pptx", "Executive briefing deck"),
            ("report.html", "Full HTML report"),
        ]
        for fname, desc in appendix_entries:
            if (self.output_dir / fname).exists():
                lines.append(f"- [`{fname}`]({fname}) — {desc}")
        lines.extend([
            "",
            "---",
            "",
            "*Report generated by NexusRecon · `master_reporter` agent · "
            "audit: `logs/audit_log.jsonl`*",
        ])

        path = self.output_dir / "master_report.md"
        path.write_text("\n".join(lines), encoding="utf-8")
        return str(path)


# ── Module-level helpers ───────────────────────────────────────────────────────

def _count_cloud_assets(cloud_intel: dict[str, Any]) -> int:
    """
    B35: count meaningful cloud assets for the executive summary.

    Previous behavior summed only S3 buckets, missing Azure tenants entirely.
    This counter includes:
      - Each cloud provider with attribution_confidence >= 0.5 (verified presence)
      - All discovered bucket/storage objects across providers
      - GCP projects discovered (when present)

    Stem-match-only entries (attribution_confidence < 0.5) are excluded so the
    operator-facing number reflects real attribution, not name-collision noise.
    """
    count = 0
    for key, data in (cloud_intel or {}).items():
        if not isinstance(data, dict):
            continue
        # Skip stem-match-only sources AND untagged sources entirely.
        # Default 0.0 (not 1.0) matches the master_report's section-4 gate so the
        # two reports agree on what counts as "verified cloud presence" — a tool
        # that doesn't emit attribution_confidence is treated as low-confidence
        # until it's explicitly tagged 1.0 by its emitter.
        if data.get("attribution_confidence", 0.0) < 0.5:
            continue
        # Count buckets / storage across all providers
        count += len(data.get("s3_buckets") or [])
        count += len(data.get("public_buckets") or [])
        count += len(data.get("gcs_buckets") or [])
        count += len(data.get("storage") or [])
        # Count verified Azure tenant as one cloud asset
        oc = data.get("openid_config") or {}
        if isinstance(oc, dict) and oc.get("found"):
            count += 1
        # Count GCP projects
        count += len(data.get("projects") or [])
    return count


def _generate_pretext_hooks(
    email: str,
    role: str,
    dept: str,
    is_breached: bool,
    is_infostealer: bool,
) -> list[str]:
    """Generate role-appropriate social engineering pretext suggestions."""
    pretexts: list[str] = []
    role_lower = role.lower()
    dept_lower = dept.lower()

    # Infostealer hit → credential reuse / account security lure
    if is_infostealer:
        pretexts.append(
            "Infostealer lure: 'We detected unusual sign-in activity on your account from a new device. "
            "Verify your identity to secure your account.' — leverages known credential exposure."
        )

    # Executive lures
    if any(kw in role_lower for kw in ["ceo", "cfo", "cto", "president", "founder"]):
        pretexts.append("Wire transfer / invoice approval — BEC-style impersonation of auditor or board member.")
        pretexts.append("Regulatory filing notification — IRS/SEC/FCA notice requiring urgent review.")

    # Finance / accounting
    if any(kw in role_lower + dept_lower for kw in ["finance", "accounting", "payroll", "cfo", "treasurer"]):
        pretexts.append("Payroll update request — 'Please review Q4 bonus structure document (attached).'")
        pretexts.append("Vendor invoice dispute — spoofed supplier domain, asks for bank detail confirmation.")

    # IT / sysadmin
    if any(kw in role_lower + dept_lower for kw in ["it", "sysadmin", "devops", "engineer", "infrastructure", "cloud"]):
        pretexts.append("Password policy enforcement — 'Your password expires in 24h, reset at [phish page].'")
        pretexts.append("AWS/Azure billing alert — unusual spend detected, click to review charges.")

    # HR / recruiting
    if any(kw in role_lower + dept_lower for kw in ["hr", "human resources", "recruiter", "talent", "people"]):
        pretexts.append("Candidate application — malicious resume attachment targeting HR file-open habit.")
        pretexts.append("Benefits enrollment reminder — 'Open enrollment closes Friday, log in to confirm.'")

    # Legal / compliance
    if any(kw in role_lower + dept_lower for kw in ["legal", "compliance", "counsel", "risk", "privacy"]):
        pretexts.append("GDPR/CCPA data subject request — spoofed regulatory body, requires portal login.")
        pretexts.append("Litigation hold notice — malicious PDF attachment.")

    # Breach history → password reset / security alert
    if is_breached and not is_infostealer:
        pretexts.append(
            "Security alert: 'Your credentials were found in a recent data breach. "
            "Reset your password now.' — highly credible given real breach exposure."
        )

    # Generic fallback
    if not pretexts:
        pretexts.append("IT helpdesk ticket — 'Your account requires re-authentication to comply with new policy.'")

    return pretexts
